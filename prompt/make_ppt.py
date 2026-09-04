"""
课设答辩 PPT 生成脚本：从项目真实的指标 JSON 自动读取数据，生成 16:9 pptx。

内容全部来自本项目实际实现（见 docs/答辩PPT大纲.md）：
  页面：封面 / 汇报框架 / 背景意义 / 方案设计 / 技术路线 / 数据构建 /
        YOLO / SVM / GA / B/S 系统 / 调试优化 / 实验结果 / 总结展望 / Q&A

数据源（backend/artifacts/）：
  yolo_train_metrics.json / yolo_test_metrics.json   YOLO 训练与测试指标
  svm_val_metrics.json                               GA 固化模型的验证集指标
  ga_best_params.json + ga_comparison.json           GA 结果与对比实验
  pipeline_image_metrics.json                        eval_pipeline 图像级指标
  ga_convergence.png / ga_comparison.png / yolo_test_confusion.png   图表

用法（venv 已激活）：
    python docs/make_ppt.py          # 输出 docs/课设答辩.pptx
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

SCRIPT_DIR = Path(__file__).resolve().parent
PROJECT_ROOT = SCRIPT_DIR.parent
ARTIFACTS = PROJECT_ROOT / "backend" / "artifacts"
OUT_PPT = SCRIPT_DIR / "课设答辩.pptx"

# 配色：深蓝主色 + 灰文字（学术答辩风格）
DARK = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
GRAY = RGBColor(0x55, 0x5A, 0x5F)
LIGHT_BG = RGBColor(0xF2, 0xF5, 0xF8)
FONT = "微软雅黑"
SLIDE_W = Inches(13.333)  # 16:9 页面宽度（add_title_bar 等工具函数无法访问 prs）


def load_json(name: str) -> dict:
    """读取 artifacts 下 JSON，缺失时返回空 dict（页面显示占位符）。"""
    path = ARTIFACTS / name
    if not path.exists():
        return {}
    return json.loads(path.read_text(encoding="utf-8"))


# 启动时读一次指标（GA 完成后运行本脚本即为最终数据）
YOLO_TRAIN = load_json("yolo_train_metrics.json")
YOLO_TEST = load_json("yolo_test_metrics.json")
SVM_TRAIN = load_json("svm_train_metrics.json")   # 分类器 AUC 存于此文件
SVM_VAL = load_json("svm_val_metrics.json")
GA_BEST = load_json("ga_best_params.json")
GA_COMP = load_json("ga_comparison.json")
PIPE = load_json("pipeline_image_metrics.json")


def fmt(v, nd=4):
    """安全格式化数字，缺失返回 '___'。"""
    if v is None:
        return "___"
    try:
        return f"{float(v):.{nd}f}"
    except (TypeError, ValueError):
        return str(v)


def svm_auc() -> str:
    """SVM 判别器验证集 AUC（svm_train_metrics.classifier_metrics.auc）。"""
    return fmt((SVM_TRAIN.get("classifier_metrics") or {}).get("auc"))


def pipe_split(split: str, key: str, nd=4) -> str:
    """取 pipeline_image_metrics.json 中某划分的图像级指标。"""
    try:
        return fmt(PIPE["splits"][split][key], nd)
    except (KeyError, TypeError):
        return "___"


# ---------------------------------------------------------------------------
# 页面工具
# ---------------------------------------------------------------------------
def blank_slide(prs: Presentation) -> object:
    """添加一张空白 16:9 幻灯片。"""
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_title_bar(slide, title: str, sub: str = "") -> None:
    """页面上方深色标题条。"""
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.45)
    tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.name = FONT
    if sub:
        p2 = tf.add_paragraph()
        p2.text = sub
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xCF, 0xDB, 0xE8)
        p2.font.name = FONT


def add_bullets(slide, items: list[str], top: float = 1.2, size: int = 18,
                gap: float = 0.5, width: float = 12.3) -> None:
    """主体要点列表：每项一行，圆点开头。"""
    box = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(width), Inches(6.2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("• " if not item.startswith("•") else "") + item
        p.font.size = Pt(size)
        p.font.name = FONT
        p.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
        p.space_after = Pt(gap * 18)
        if item.startswith("→"):
            p.level = 1
            p.font.size = Pt(size - 2)
            p.font.color.rgb = GRAY
    return box


def add_note(slide, text: str) -> None:
    """页面底部灰色小字（页面提示/出处，不放讲稿内容）。"""
    box = slide.shapes.add_textbox(Inches(0.55), Inches(6.85), Inches(12.3), Inches(0.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.name = FONT
    p.font.color.rgb = GRAY


def add_image(slide, path: str, left: float, top: float, width: float) -> None:
    """插入图片（按宽度等比缩放）；文件缺失时跳过。"""
    p = Path(path)
    if not p.exists():
        return
    slide.shapes.add_picture(str(p), Inches(left), Inches(top), width=Inches(width))


def add_speaker_note(slide, text: str) -> None:
    """写入演讲者备注（打印/放映时可见，便于脱稿）。"""
    slide.notes_slide.notes_text_frame.text = text


# ---------------------------------------------------------------------------
# 页面构建
# ---------------------------------------------------------------------------
def build(prs: Presentation) -> None:
    # P1 封面
    s = blank_slide(prs)
    bg = s.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK
    bg.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(2.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "基于深度学习的轴承表面缺陷智能检测系统"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.name = FONT
    p2 = tf.add_paragraph()
    p2.text = "—— YOLOv8 一级定位 + SVM 二级判别 + 遗传算法超参数优化"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(0xCF, 0xDB, 0xE8)
    p2.font.name = FONT
    tb2 = s.shapes.add_textbox(Inches(0.8), Inches(4.9), Inches(8.0), Inches(1.2))
    tf2 = tb2.text_frame
    for line in ("《制造智能技术》课程设计答辩", "姓名：________    学号：________", "指导教师：________"):
        para = tf2.paragraphs[0] if line == "《制造智能技术》课程设计答辩" else tf2.add_paragraph()
        para.text = line
        para.font.size = Pt(17)
        para.font.name = FONT
        para.font.color.rgb = RGBColor(0xCF, 0xDB, 0xE8)
    add_speaker_note(s, "各位老师好，我的课设题目是《基于深度学习的轴承表面缺陷智能检测系统》。")

    # P2 汇报框架
    s = blank_slide(prs)
    add_title_bar(s, "汇报框架")
    add_bullets(s, [
        "① 背景调研与研究意义",
        "② 方案设计：两级检测架构",
        "③ 数据资源构建",
        "④ 详细开发内容（三个技术方向 + B/S 系统）",
        "⑤ 调试与问题解决",
        "⑥ 实验结果 · 总结与展望",
    ], top=1.4, size=20, gap=0.62)
    add_speaker_note(s, "下面按这个顺序汇报，约三分钟。")

    # P3 背景与意义
    s = blank_slide(prs)
    add_title_bar(s, "研究背景与意义")
    add_bullets(s, [
        "轴承表面缺陷（划痕等）直接影响寿命与可靠性，是失效的重要诱因",
        "传统人工目检：劳动强度大、效率低、标准不统一、易漏检",
        "深度学习：自动学习缺陷特征，检测速度与一致性远超人工，适合产线重复性视觉任务",
        "本课设：公开数据集 KolektorSDD + 三个技术方向（计算机视觉 / 机器学习 / 智能优化）",
    ], top=1.3, size=19, gap=0.6)
    add_speaker_note(s, "轴承是机械装备的核心零件。产线目检存在这些问题，深度学习适合这类任务。我把三个技术方向组合成一套系统。")

    # P4 方案设计（核心页）
    s = blank_slide(prs)
    add_title_bar(s, "方案设计：为什么是两级架构（核心决策）")
    add_bullets(s, [
        "实测：小数据集训练的 YOLO 置信度标定差——候选框最高置信度仅 0.013，单一阈值无法兼顾漏检与误报",
        "→ 一级 YOLOv8：低阈值高召回，只负责“找出所有可疑区域”（宁可多检不漏）",
        "→ 二级 SVM：HOG+LBP+统计特征（1826 维）判别，过滤误检",
        "→ 遗传算法：联合优化 SVM 超参数 + 两级阈值，适应度 = 验证集流水线 F1",
        f"两级流水线把图像级 F1 从 0.28 提升到 {pipe_split('test', 'f1')}（test）",
    ], top=1.25, size=18, gap=0.52)
    add_speaker_note(s, "核心决策。如果只靠 YOLO 一个模型，实测置信度整体偏低，单一阈值怎么设都不好。所以拆成两级：第一级宁可多检不漏，第二级精细判别，再用 GA 把阈值和分类器参数一起优化。")

    # P5 技术路线（架构图）
    s = blank_slide(prs)
    add_title_bar(s, "技术路线：离线优化 + 在线推理")
    add_bullets(s, [
        "离线阶段：数据构建 → YOLOv8 训练 → 候选框特征提取 → SVM 训练 → GA 超参数优化 → 固化最优模型",
        "在线阶段：Vue3 上传图像 → FastAPI 两级流水线推理 → PostgreSQL 落库 / Redis 缓存",
    ], top=1.15, size=17, gap=0.45)
    add_speaker_note(s, "整体分两段：离线把模型和参数训练优化好；在线是一个完整的浏览器/服务器系统。")

    # P6 数据构建
    s = blank_slide(prs)
    add_title_bar(s, "数据资源构建（KolektorSDD）")
    add_bullets(s, [
        "399 张工业表面缺陷图像：52 张有缺陷 / 347 张正常，附像素级掩码（无检测框标注）",
        "掩码 → 二值化(>127) → 连通域分析 → 外接矩形，自动生成目标检测标注（min_area=25 滤噪）",
        "划分：按“物理件”分组——同一零件的多张子图必须同属一个划分，防止数据泄漏",
    ], top=1.3, size=18, gap=0.55)
    add_speaker_note(s, "数据集只有掩码没有检测框，我写脚本转换。划分特别注意：同一零件的多个视角必须整组划分。")

    # P7 YOLOv8
    s = blank_slide(prs)
    add_title_bar(s, "开发内容 ①：YOLOv8 一级定位（计算机视觉）")
    add_bullets(s, [
        "模型 yolov8n，imgsz=640，100 epochs，数据增强（翻转/色彩/马赛克）应对小数据集",
        f"验证集 mAP50 ≈ {fmt(YOLO_TRAIN.get('mAP50'))}；官方 Test 集 Recall ≈ {fmt(YOLO_TEST.get('recall'))}",
        "定位策略：候选阈值压低到 0.001 最大化召回，每图取置信度 top-50 候选送二级判别",
    ], top=1.3, size=18, gap=0.55)
    if (ARTIFACTS / "yolo_test_confusion.png").exists():
        add_image(s, str(ARTIFACTS / "yolo_test_confusion.png"), 8.3, 3.3, 4.0)
        add_note(s, "右图：YOLOv8 官方 Test 集混淆矩阵")
    add_speaker_note(s, "一级定位 YOLOv8n。召回率九成说明找可疑区域这个任务完成得不错，精度交给二级。")

    # P8 SVM
    s = blank_slide(prs)
    add_title_bar(s, "开发内容 ②：SVM 二级判别（机器学习）")
    add_bullets(s, [
        "候选框统一缩放 64×64 → 特征：HOG(1764 维，手写实现) + LBP 旋转不变均匀模式(59 维) + 统计特征(3 维)",
        "正样本：候选框与真值框 IoU ≥ 0.5；负样本：YOLO 误检 + 正常件随机块 + 边缘条带块",
        f"验证集分类 AUC ≈ {svm_auc()}",
    ], top=1.3, size=18, gap=0.55)
    add_speaker_note(s, "SVM 判别。HOG 描述边缘方向、LBP 描述纹理。HOG 是我自己用 numpy 手写的，原理可以讲清楚。负样本三类来源，解决正常件假阳性。")

    # P9 GA
    s = blank_slide(prs)
    add_title_bar(s, "开发内容 ③：遗传算法超参数优化（智能优化）")
    if GA_BEST:
        ga_line = f"GA 最优适应度 F1 = {fmt(GA_BEST.get('fitness'))}（对比实验见实验结果页）"
    else:
        ga_line = "实验设置：种群 16 × 20 代 + 随机搜索 40 组 + 网格搜索 24 组，三方法对比（结果见 artifacts/）"
    add_bullets(s, [
        "基因（5 维混合编码）：logC、logγ、核函数{线性/多项式/RBF}、YOLO 置信度阈值、SVM 概率阈值",
        "算子：锦标赛选择(k=3) · 单点交叉(Pc=0.8) · 混合变异(Pm=0.1) · 精英保留 2 名",
        "适应度 = 验证集两级流水线 F1；多进程并行评估加速（12 进程）",
        ga_line,
    ], top=1.25, size=18, gap=0.52)
    add_speaker_note(s, "手写 GA。把 SVM 两个超参数和流水线两个阈值当基因联合优化，适应度直接用验证集 F1。做了随机搜索和网格搜索对比。")

    # P10 B/S 系统
    s = blank_slide(prs)
    add_title_bar(s, "开发内容 ④：B/S 应用系统")
    add_bullets(s, [
        "FastAPI 后端：两级流水线推理（懒加载单例）→ 检测记录落库 → 结果缓存回填",
        "Vue3 + Element Plus + ECharts 前端四页面：在线检测 / 检测历史 / 统计概览 / 模型信息",
        "PostgreSQL 存储检测记录与模型元数据；Redis 按图像 MD5 缓存 600s（重复上传免推理）",
        "单元测试 17 个（pytest，SQLite 内存库隔离）+ docker-compose 一键部署配置",
    ], top=1.25, size=18, gap=0.52)
    add_speaker_note(s, "把模型包装成可用系统。后端 FastAPI，前端四页面，PG 存记录、Redis 缓存，附带 Docker 部署。")

    # P11 调试
    s = blank_slide(prs)
    add_title_bar(s, "调试与问题解决（摘选）")
    add_bullets(s, [
        "Windows 中文路径：OpenCV 读中文路径失败 → 封装 np.fromfile + imdecode 统一解决",
        "置信度标定差：诊断脚本实测最高 conf 仅 0.013 → 阈值降 0.001 + 每图 top-50 候选送 SVM",
        "正常件假阳性：误报 7 框 → 负样本多样性增强（随机块 + 边缘条带块）→ GA 优化阈值进一步过滤",
        "统计时区坑：SQLite CURRENT_TIMESTAMP 为 UTC，近 7 天统计漏当天 → 应用层本地时间",
    ], top=1.25, size=18, gap=0.52)
    add_speaker_note(s, "踩了四个有代表性的坑，排查过程都记录在规格文档里。")

    # P12 实验结果
    s = blank_slide(prs)
    add_title_bar(s, "实验结果")
    if GA_COMP:
        ga_cmp_line = (f"优化方法对比（验证集 F1）：GA {fmt(GA_COMP.get('ga', {}).get('best_fitness'))}"
                       f" / 随机 {fmt(GA_COMP.get('random', {}).get('best_fitness'))}"
                       f" / 网格 {fmt(GA_COMP.get('grid', {}).get('best_fitness'))}")
        cmp_note = "下排左：GA 逐代收敛曲线；下排右：三种优化方法最优 F1 对比"
        cmp_speech = "最终指标。GA 收敛曲线说明种群在十几代内收敛，对比图显示 GA 优于随机搜索和网格搜索。"
        cmp_img = True
    else:
        ga_cmp_line = "GA / 随机搜索 / 网格搜索对比实验已设计，运行完成后回填（重新运行 docs/make_ppt.py 自动更新）"
        cmp_note = ""
        cmp_speech = "最终指标。GA 对比实验在离线脚本中运行，结果可在 artifacts 目录查看。"
        cmp_img = False
    add_bullets(s, [
        f"YOLOv8：验证集 mAP50 {fmt(YOLO_TRAIN.get('mAP50'))} / 官方 Test 集 Recall {fmt(YOLO_TEST.get('recall'))}（高召回定位）",
        f"SVM 判别器：验证集 AUC ≈ {svm_auc()}",
        f"两级流水线（图像级，top-50 候选）：val F1 {pipe_split('val', 'f1')} · test F1 {pipe_split('test', 'f1')}（P {pipe_split('test', 'precision')} / R {pipe_split('test', 'recall')} / ACC {pipe_split('test', 'accuracy')}）",
        ga_cmp_line,
    ], top=1.2, size=16, gap=0.42)
    if cmp_img:
        add_image(s, str(ARTIFACTS / "ga_convergence.png"), 0.6, 4.35, 5.7)
        add_image(s, str(ARTIFACTS / "ga_comparison.png"), 6.7, 4.35, 5.7)
        add_note(s, cmp_note)
    add_speaker_note(s, cmp_speech)

    # P13 总结展望
    s = blank_slide(prs)
    add_title_bar(s, "总结与展望")
    add_bullets(s, [
        "完成：三算法各司其职的两级检测系统 + 完整 B/S 应用 + 规格/测试/Docker 全套工程交付",
        "关键收获：小数据集下“分工架构 + 联合优化”优于单模型调阈值",
        "不足：缺陷样本仅 52 张；YOLO 置信度标定差只能工程补救",
        "展望：采集真实轴承数据 / 升级检测模型（二级判别架构不变）/ 对接工业相机产线",
    ], top=1.3, size=18, gap=0.55)
    add_speaker_note(s, "总结三句话：架构上两级分工解决了小数据集问题；工程上整合成完整系统；未来换数据换模型，架构不用改。")

    # P14 Q&A
    s = blank_slide(prs)
    bg = s.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK
    bg.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "感谢聆听，请各位老师批评指正"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER


def main() -> None:
    """生成 PPT。"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    build(prs)
    prs.save(OUT_PPT)
    print(f"PPT 已生成: {OUT_PPT}")
    print(f"  幻灯片 {len(prs.slides)} 页；指标读取自 backend/artifacts/（缺失显示 '___'）")


if __name__ == "__main__":
    sys.stdout.reconfigure(encoding="utf-8")
    main()
